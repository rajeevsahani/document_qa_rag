import io
import os
from docx import Document as DocxDocument
import PyPDF2
from pptx import Presentation
from fastapi import HTTPException
from embeddings import Embeddings
from repository import DocumentRepository
from llm import LLMModel

class DocumentService:
    def __init__(self, repository: DocumentRepository):
        self.repository = repository
        self.embedder = Embeddings()
        self.llm = LLMModel()  # Initialize LLM model

    async def process_file(self, file):
        file_name = file.filename
        file_extension = os.path.splitext(file_name)[1].lower()
        contents = await file.read()
        text_content = ""

        if file_extension == ".txt":
            text_content = contents.decode("utf-8")
        elif file_extension == ".docx":
            doc = DocxDocument(io.BytesIO(contents))
            text_content = "\n".join([para.text for para in doc.paragraphs])
        elif file_extension == ".pdf":
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
            text_content = "".join([page.extract_text() for page in pdf_reader.pages])
        elif file_extension in [".ppt", ".pptx"]:
            presentation = Presentation(io.BytesIO(contents))
            text_content = "\n".join(
                [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")]
            )
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        embedding = self.embedder.get_embeddings(text_content)
        file_id = await self.repository.add_document(file_name, text_content, embedding)
        return {"message": "File uploaded successfully", "file_id": str(file_id)}

    async def list_documents(self):
        return await self.repository.get_documents()

    # async def search_documents(self, query: str):
    #     query_embedding = self.embedder.get_embeddings(query).tolist() 
    #     return await self.repository.search_similar_documents(query, query_embedding)

    async def search_documents(self, query: str):
        query_embedding = self.embedder.get_embeddings(query).tolist() 
        retrieved_docs = await self.repository.search_similar_documents(query, query_embedding)
        
        # Combine retrieved document contents for LLM input
        context = "\n\n".join([doc["content"] for doc in retrieved_docs["results"]])

        # Generate an answer using an LLM model
        response = self.llm.generate_answer(query, context)

        return {
            "query": query,
            "answer": response,
            "retrieved_documents": retrieved_docs["results"]
        }
